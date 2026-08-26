import re
import json
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from app.database import Base
from app.models import CasePackage, CaseTask, User
from app.seed import _recover_interrupted_generation_tasks
from app.services.orchestrator import _task_context, _validate_agent_output
from app.services.package_builder import (
    build_structured_package,
    build_teaching_flow,
    ensure_mock_body_length,
    fit_teaching_flow_to_class_hours,
    normalize_case_package,
)
from app.services.grounded_case_service import generation_preflight_error
from app.services.video_service import search_official_videos
from app.services.objective_generator import generate_objectives
from app.services.pptx_export_service import build_ppt_outline, export_pptx
from app.services.skill_loader import validate_package_with_skill


def make_task(class_hours: int = 1) -> CaseTask:
    return CaseTask(
        title="系统工程——综合集成方法论",
        subject="管理学",
        course_name="系统工程",
        case_type="决策型",
        difficulty="中级",
        target_audience="本科",
        target_words=2800,
        learning_objectives=[
            "分析综合集成方法的适用条件",
            "评价不同利益相关者的系统边界判断",
            "设计可执行的综合集成方案",
        ],
        workflow_template="sequential_standard",
        config={"class_hours": class_hours, "special_requirements": ""},
        status="draft",
    )


def flow_minutes(flow: str) -> int:
    return sum(int(value) for value in re.findall(r"(\d+)\s*(?:min|分钟)", flow, flags=re.I))


class GenerationReliabilityTests(unittest.TestCase):
    def test_objective_brief_turns_teacher_intent_into_observable_goals(self) -> None:
        result = generate_objectives({
            "title": "制造企业数字化转型中的组织阻力",
            "subject": "管理学",
            "course_name": "组织行为学",
            "case_type": "问题诊断",
            "difficulty": "中级",
            "target_audience": "本科",
            "variant": 0,
            "objective_brief": {
                "learning_challenge": "学生会背理论但不会用证据解释组织阻力",
                "desired_performance": "比较两种干预方案并作出有依据的选择",
                "required_concepts": "组织变革阻力与系统反馈",
                "assessment_evidence": "一张因果关系图和三分钟小组答辩",
            },
        })
        self.assertEqual(result["quality_score"], 100)
        self.assertEqual(len(result["objectives"]), 3)
        self.assertIn("学生会背理论", result["objectives"][0])
        self.assertIn("组织变革阻力与系统反馈", result["objectives"][1])
        self.assertIn("因果关系图", result["objectives"][2])
        self.assertTrue(all(item["passed"] for item in result["quality_checks"]))

    def test_strict_teacher_brief_uses_reviewed_sources_and_never_fictional_template(self) -> None:
        task = make_task(1)
        task.learning_objectives = [
            "以钱学森为代表的综合集成方法论，匹配适合和前沿的案例",
            "通过闭环和反馈思想解决实际问题，培养本科生系统工程思维",
            "包含课程思政和灵魂人物",
        ]
        task.config = {"class_hours": 1, "special_requirements": "生成的案例不能虚假，数据准确，切忌编造。"}
        package = build_structured_package(task)
        visible = "".join(str(package["body"].get(key, "")) for key in ("background", "narrative", "decision_point"))
        self.assertEqual(package["meta"]["content_mode"], "source_grounded")
        self.assertEqual(len(package["evidence_sources"]), 10)
        self.assertEqual(len(package["visual_assets"]), 10)
        self.assertEqual(len(package["video_resources"]), 10)
        json.dumps(package, ensure_ascii=False)
        self.assertNotIn("某组织", visible)
        self.assertNotIn("陈启明", visible)
        self.assertTrue(2660 <= len(re.sub(r"\s+", "", visible)) <= 2940)
        validation = validate_package_with_skill(package, _task_context(task))
        self.assertTrue(validation["passed"], validation["issues"])

    def test_video_recommendations_are_course_scoped_and_traceable(self) -> None:
        videos = search_official_videos("系统工程 综合集成方法论 钱学森 中国载人航天", 10)
        self.assertEqual(len(videos), 10)
        for video in videos:
            self.assertTrue(video["source_page_url"].startswith("https://"))
            self.assertTrue(video["video_url"].startswith("https://"))
            self.assertEqual(video["trust_level"], "official")
            self.assertTrue(video["usage"])
        self.assertEqual(search_official_videos("完全无关的古典音乐赏析课程", 10), [])

    def test_strict_factual_case_without_reviewed_profile_stops_before_fabrication(self) -> None:
        task = make_task(1)
        task.title = "尚未建立事实资料包的新课程"
        task.config = {"class_hours": 1, "special_requirements": "真实案例，不得虚构。"}
        self.assertIsNotNone(generation_preflight_error(task))

    def test_all_supported_class_hours_fit_quality_gate(self) -> None:
        for hours in range(1, 9):
            with self.subTest(hours=hours):
                task = make_task(hours)
                package = build_structured_package(task)
                ensure_mock_body_length(package, task)
                flow = package["instructor_guide"]["teaching_flow"]
                self.assertEqual(flow_minutes(flow), hours * 45)
                validation = validate_package_with_skill(package, _task_context(task))
                timing_errors = [
                    issue for issue in validation["issues"]
                    if issue.get("code") in {"timing_overflow", "timing_missing"}
                ]
                self.assertEqual(timing_errors, [])

    def test_ppt_uses_editorial_galleries_and_teacher_method_structure(self) -> None:
        task = make_task(2)
        task.config["special_requirements"] = "以钱学森综合集成方法论为知识点，必须使用真实案例并融入课程思政。"
        package = build_structured_package(task)
        outline = build_ppt_outline(package, {"density": "standard", "audience": "teacher"})
        kinds = [slide["kind"] for slide in outline["slides"]]
        self.assertIn("visual_gallery", kinds)
        self.assertNotIn("visual", kinds)
        self.assertIn("method_map", kinds)
        self.assertIn("ideology", kinds)
        galleries = [slide for slide in outline["slides"] if slide["kind"] == "visual_gallery"]
        self.assertTrue(all(1 <= len(slide["items"]) <= 3 for slide in galleries))
        self.assertEqual(sum(len(slide["items"]) for slide in galleries), 10)
        self.assertLessEqual(sum(kind == "sources" for kind in kinds), 2)
        self.assertLessEqual(sum(kind == "videos" for kind in kinds), 2)
        self.assertGreaterEqual(outline["design_metrics"]["visual_pages"], 4)

    def test_redesigned_ppt_is_a_valid_editable_widescreen_file(self) -> None:
        package = build_structured_package(make_task(2))
        from app.services import pptx_export_service
        with TemporaryDirectory() as directory:
            original_export_dir = pptx_export_service.settings.export_dir
            try:
                pptx_export_service.settings.export_dir = directory
                path = Path(export_pptx(package, package["meta"]["title"], version=8))
                deck = Presentation(str(path))
                self.assertGreater(len(deck.slides), 20)
                self.assertEqual(round(deck.slide_width / deck.slide_height, 2), 1.78)
                self.assertTrue(any(shape.has_text_frame for slide in deck.slides for shape in slide.shapes))
            finally:
                pptx_export_service.settings.export_dir = original_export_dir

    def test_overflowing_model_schedule_is_repaired(self) -> None:
        package = {"instructor_guide": {"teaching_flow": "阅读(20min)→讨论(25min)→汇报(20min)"}}
        self.assertTrue(fit_teaching_flow_to_class_hours(package, 1))
        self.assertEqual(flow_minutes(package["instructor_guide"]["teaching_flow"]), 45)
        self.assertFalse(fit_teaching_flow_to_class_hours(package, 1))

    def test_internal_notes_are_removed_from_every_student_body_field(self) -> None:
        package = {
            "body": {
                "background": "背景正文。【学科注释】共同注释",
                "narrative": "叙事正文。【学科注释】共同注释",
                "decision_point": "决策正文。【学科注释】另一注释",
            }
        }
        normalize_case_package(package)
        visible = "".join(package["body"].values())
        self.assertNotIn("【学科注释】", visible)
        self.assertEqual(package["domain_context"]["notes"], "共同注释\n另一注释")

    def test_pedagogy_contract_rejects_incomplete_release_fields(self) -> None:
        task = make_task(1)
        errors = _validate_agent_output(
            "PedagogyDesigner",
            {
                "discussion_questions": [
                    {"level": level, "question": f"问题{i}", "teaching_intent": "意图"}
                    for i, level in enumerate(["理解", "分析", "评价", "创造", "分析"], 1)
                ],
                "instructor_guide": {"teaching_flow": build_teaching_flow(1), "key_points": []},
                "alignment_matrix": [
                    {"objective_id": f"LO{i}", "case_section": "正文", "activity": "讨论", "assessment": ""}
                    for i in range(1, 4)
                ],
            },
            task,
        )
        self.assertTrue(any("key_points" in error for error in errors))
        self.assertTrue(any("common_misconceptions" in error for error in errors))
        self.assertTrue(any("字段不完整" in error for error in errors))

    def test_restart_recovers_running_tasks_without_leaving_them_stuck(self) -> None:
        engine = create_engine("sqlite:///:memory:")
        Base.metadata.create_all(engine)
        session = sessionmaker(bind=engine)()
        try:
            user = User(email="recovery@example.com", password_hash="hash", name="测试")
            session.add(user)
            session.flush()
            with_package = make_task(1)
            with_package.user_id = user.id
            with_package.status = "running"
            without_package = make_task(1)
            without_package.user_id = user.id
            without_package.title = "被重启中断的新任务"
            without_package.status = "running"
            session.add_all([with_package, without_package])
            session.flush()
            session.add(CasePackage(task_id=with_package.id, version=1, package={}, status="finalized"))
            session.commit()

            _recover_interrupted_generation_tasks(session)
            session.commit()

            self.assertEqual(with_package.status, "finalized")
            self.assertIsNone(with_package.error_message)
            self.assertEqual(without_package.status, "failed")
            self.assertIn("重新生成", without_package.error_message)
        finally:
            session.close()


if __name__ == "__main__":
    unittest.main()
