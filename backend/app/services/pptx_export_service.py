"""Build editable, classroom-ready PPTX decks from a structured teaching case package."""

from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.config import settings
from app.services.material_service import get_cached_material_image, resolve_package_materials


THEMES = {
    "academic": {
        "style": "academic",
        "primary": "312E81", "accent": "4F46E5", "soft": "EEF2FF",
        "background": "F8FAFC", "surface": "FFFFFF", "text": "172033", "muted": "64748B",
        "accent2": "06B6D4", "success": "059669", "warning": "D97706", "danger": "DC2626",
    },
    "modern": {
        "style": "modern",
        "primary": "0F3D56", "accent": "0D9488", "soft": "CCFBF1",
        "background": "F6FAFA", "surface": "FFFFFF", "text": "132A35", "muted": "607580",
        "accent2": "0284C7", "success": "16A34A", "warning": "EA580C", "danger": "DC2626",
    },
    "minimal": {
        "style": "editorial",
        "primary": "111827", "accent": "D97706", "soft": "FEF3C7",
        "background": "FAFAF9", "surface": "FFFFFF", "text": "1C1917", "muted": "78716C",
        "accent2": "0F766E", "success": "15803D", "warning": "D97706", "danger": "B91C1C",
    },
    "executive": {
        "style": "executive",
        "primary": "0B1F3A", "accent": "C08A2D", "soft": "F4E8CE",
        "background": "F5F7FA", "surface": "FFFFFF", "text": "162338", "muted": "667085",
        "accent2": "2E6F95", "success": "16866A", "warning": "C76B20", "danger": "C23B4A",
    },
    "vibrant": {
        "style": "vibrant",
        "primary": "3B185F", "accent": "7C3AED", "soft": "F3E8FF",
        "background": "FFF9FC", "surface": "FFFFFF", "text": "24132F", "muted": "75657F",
        "accent2": "E85D75", "success": "079669", "warning": "E07A1F", "danger": "D9365B",
    },
}

THEME_NAMES = {
    "academic": "学术靛蓝",
    "modern": "现代青绿",
    "minimal": "极简暖灰",
    "executive": "高管深蓝",
    "vibrant": "活力紫红",
}

DENSITY = {
    "concise": {"bullets": 3, "bullet_chars": 52, "narrative_slides": 3, "questions": 3, "gallery_size": 4},
    "standard": {"bullets": 4, "bullet_chars": 68, "narrative_slides": 5, "questions": 3, "gallery_size": 3},
    "detailed": {"bullets": 5, "bullet_chars": 88, "narrative_slides": 8, "questions": 2, "gallery_size": 3},
}


def _clean_text(value: Any) -> str:
    text = re.sub(r"【(?:背景|叙事|案例正文)】", "", str(value or ""))
    return re.sub(r"\s+", " ", text).strip()


def _split_sentences(text: str) -> list[str]:
    normalized = str(text or "").replace("\r\n", "\n").strip()
    parts = re.split(r"(?<=[。！？；])\s*|\n+", normalized)
    return [
        cleaned for part in parts
        if (cleaned := _clean_text(part)) and re.sub(r"[\W_]+", "", cleaned)
    ]


def _fit_segment(sentence: str, max_chars: int) -> list[str]:
    if len(sentence) <= max_chars:
        return [sentence]
    clauses = [part.strip() for part in re.split(r"(?<=[，、：])", sentence) if part.strip()]
    result: list[str] = []
    current = ""
    for clause in clauses:
        if current and len(current) + len(clause) > max_chars:
            result.append(current)
            current = clause
        else:
            current += clause
    if current:
        result.append(current)
    hard_limit = int(max_chars * 1.3)
    return [item if len(item) <= hard_limit else item[:hard_limit] + "…" for item in result]


def _content_pages(text: str, density: str, max_slides: int | None = None) -> list[list[str]]:
    rule = DENSITY[density]
    segments: list[str] = []
    for sentence in _split_sentences(text):
        segments.extend(_fit_segment(sentence, rule["bullet_chars"]))
    pages = [segments[index : index + rule["bullets"]] for index in range(0, len(segments), rule["bullets"])]
    if max_slides and len(pages) > max_slides:
        kept = pages[:max_slides]
        if sum(len(page) for page in pages[max_slides:]):
            kept[-1][-1] = kept[-1][-1].rstrip("。…") + "……（其余细节请结合案例原文阅读）"
        return kept
    return pages or [["暂无内容"]]


def _paginate_items(
    items: list[Any],
    *,
    max_items: int,
    max_chars: int,
    text_getter,
) -> list[list[Any]]:
    """按项目数和文字量双重分页，杜绝固定行高硬塞长内容。"""
    pages: list[list[Any]] = []
    current: list[Any] = []
    current_chars = 0
    for item in items:
        weight = max(len(_clean_text(text_getter(item))), 1)
        if current and (len(current) >= max_items or current_chars + weight > max_chars):
            pages.append(current)
            current = []
            current_chars = 0
        current.append(item)
        current_chars += weight
    if current:
        pages.append(current)
    return pages or [[]]


def _decision_pages(text: str, density: str) -> list[list[str]]:
    segments: list[str] = []
    segment_limit = {"concise": 58, "standard": 68, "detailed": 76}[density]
    for sentence in _split_sentences(text):
        segments.extend(_fit_segment(sentence, segment_limit))
    return _paginate_items(
        segments,
        max_items=2,
        max_chars={"concise": 112, "standard": 128, "detailed": 145}[density],
        text_getter=lambda item: item,
    )


def _speaker_note(title: str, purpose: str, prompts: list[str] | None = None, minutes: int = 3) -> str:
    lines = [f"【本页目的】{purpose}", f"【建议时长】{minutes} 分钟"]
    if prompts:
        lines.append("【讲解提示】")
        lines.extend(f"- {item}" for item in prompts if item)
    lines.append(f"【转场】完成“{title}”后，引导学生进入下一环节。")
    return "\n".join(lines)


def _section_slide(title: str, number: str, kicker: str) -> dict[str, Any]:
    return {
        "kind": "section", "title": title, "number": number, "kicker": kicker,
        "notes": _speaker_note(title, "提示课堂进入新的学习阶段。", [kicker], 1),
    }


def _visual_gallery_slides(materials: list[dict[str, Any]], density: str) -> list[dict[str, Any]]:
    """Turn evidence images into editorial spreads instead of one repetitive slide per image."""
    page_size = int(DENSITY[density]["gallery_size"])
    pages: list[dict[str, Any]] = []
    for index in range(0, len(materials), page_size):
        page_items = materials[index : index + page_size]
        pages.append({
            "kind": "visual_gallery", "section": "案例证据",
            "title": "真实现场：从图像中读出系统" if index == 0 else "真实现场：证据切片",
            "items": page_items,
            "notes": _speaker_note(
                "视觉证据", "把图片作为观察和验证材料，而不是装饰。",
                [*[item.get("caption", "") for item in page_items], "追问：图片能证明什么，又不能证明什么？"], 3,
            ),
        })
    return pages


def _design_metrics(slides: list[dict[str, Any]]) -> dict[str, Any]:
    visual_pages = sum(1 for slide in slides if slide["kind"] in {"cover", "visual_gallery"})
    activity_pages = sum(1 for slide in slides if slide["kind"] in {"decision", "questions", "flow", "alignment"})
    appendix_pages = sum(1 for slide in slides if slide["kind"] in {"sources", "videos"})
    return {
        "visual_pages": visual_pages,
        "activity_pages": activity_pages,
        "appendix_pages": appendix_pages,
        "quality_label": "自适应课堂排版",
        "quality_summary": "长内容会按可读字量自动续页，并在导出前执行文字密度与画布边界检查。",
    }


def build_ppt_outline(package: dict[str, Any], options: dict[str, Any] | None = None) -> dict[str, Any]:
    options = options or {}
    density = options.get("density", "standard")
    if density not in DENSITY:
        density = "standard"
    audience = options.get("audience", "teacher")
    if audience not in {"student", "teacher"}:
        audience = "teacher"
    theme = options.get("theme", "executive")
    if theme not in THEMES:
        theme = "executive"
    mode = options.get("mode", "lecture")
    if mode not in {"lecture", "workshop", "visual"}:
        mode = "lecture"
    include_speaker_notes = bool(options.get("include_speaker_notes", True))

    meta = package.get("meta") or {}
    body = package.get("body") or {}
    title = _clean_text(meta.get("title") or "教学案例")
    course = _clean_text(meta.get("course") or meta.get("subject") or "课程教学")
    materials = resolve_package_materials(package, limit=10)
    teacher_requirements = package.get("teacher_requirements") or {}
    ideology = package.get("course_ideology") or {}
    slides: list[dict[str, Any]] = [
        {
            "kind": "cover", "title": title,
            "subtitle": f"{course} · {meta.get('case_type', '教学案例')} · {meta.get('target_audience', '')}",
            "asset_id": materials[0]["id"] if materials else None,
            "eyebrow": f"{meta.get('subject', '课程案例')} / {meta.get('difficulty', '标准难度')}",
            "notes": _speaker_note(title, "建立案例主题与课程期待。", ["不要提前透露决策结论。", "用一句现实问题引发学生好奇。"], 2),
        }
    ]

    objectives = [
        _clean_text(item.get("description"))
        for item in package.get("learning_objectives") or []
        if _clean_text(item.get("description"))
    ]
    questions = package.get("discussion_questions") or []
    guide = package.get("instructor_guide") or {}
    slides.append({
        "kind": "agenda", "title": "今天如何推进这个案例",
        "items": ["进入情境", "识别角色", "追踪冲突", "做出决策", "复盘迁移"],
        "notes": _speaker_note("课堂路径", "向学生说明本课不是寻找唯一答案，而是练习有依据的判断。", ["说明课堂参与方式。"], 2),
    })
    slides.append({
        "kind": "overview", "section": "课程导入", "title": "案例一览",
        "items": [
            {"label": "课程", "value": course},
            {"label": "对象", "value": _clean_text(meta.get("target_audience") or "学习者")},
            {"label": "类型", "value": _clean_text(meta.get("case_type") or "教学案例")},
            {"label": "难度", "value": _clean_text(meta.get("difficulty") or "标准")},
        ],
        "stats": [
            {"value": len(objectives), "label": "学习目标"},
            {"value": len(body.get("characters") or []), "label": "关键角色"},
            {"value": len(questions), "label": "研讨问题"},
        ],
        "notes": _speaker_note("案例一览", "快速建立学习边界和任务规模。", ["确认学生具备必要先修知识。"], 2),
    })
    if objectives:
        objective_pages = _paginate_items(
            objectives,
            max_items=2,
            max_chars=260,
            text_getter=lambda item: item,
        )
        objective_levels = [_clean_text(item.get("level")) for item in package.get("learning_objectives") or []]
        objective_offset = 0
        for page_index, page_items in enumerate(objective_pages, 1):
            page_levels = objective_levels[objective_offset : objective_offset + len(page_items)]
            objective_offset += len(page_items)
            slides.append({
                "kind": "objectives", "section": "课程导入",
                "title": "从识别到决策：本课学习目标" if page_index == 1 else f"本课学习目标 · {page_index}",
                "items": page_items,
                "levels": page_levels,
                "start_number": objective_offset - len(page_items) + 1,
                "notes": _speaker_note("学习目标", "让学生知道课堂结束时需要交付什么能力。", page_items, 3),
            })

    knowledge_anchors = [str(item) for item in teacher_requirements.get("knowledge_anchors") or [] if str(item).strip()]
    if knowledge_anchors:
        slides.append({
            "kind": "method_map", "section": "课程导入", "title": "方法不是名词：它如何形成认识闭环",
            "items": knowledge_anchors[:6],
            "problem": _clean_text(teacher_requirements.get("teaching_problem")),
            "notes": _speaker_note("方法路径", "把教师指定的知识点组织为可执行的认识过程。", knowledge_anchors[:6], 4),
        })

    if ideology.get("figure") or ideology.get("themes"):
        slides.append({
            "kind": "ideology", "section": "课程导入", "title": "价值坐标：方法背后的责任与选择",
            "figure": _clean_text(ideology.get("figure")),
            "themes": [_clean_text(item) for item in ideology.get("themes") or []],
            "implementation": _clean_text(ideology.get("implementation")),
            "notes": _speaker_note("价值坐标", "用可核验人物与实践呈现价值观，不虚构名言或煽情故事。", ideology.get("themes") or [], 3),
        })

    slides.append(_section_slide("进入案例现场", "01", "先理解情境，再急于判断"))

    for index, items in enumerate(_content_pages(body.get("background", ""), density), 1):
        slides.append({
            "kind": "story", "section": "案例情境",
            "title": "案例背景" if index == 1 else f"案例背景 · {index}", "items": items,
            "chapter": "CONTEXT", "emphasis": index % 2, "visual_mode": mode == "visual",
            "notes": _speaker_note("案例背景", "交代决策发生的环境、目标和约束。", [items[0] if items else ""], 3),
        })

    # Front-load only the first evidence spread; the rest appears after the story to vary pacing.
    visual_pages = _visual_gallery_slides(materials, density)
    if visual_pages:
        slides.append(visual_pages[0])

    characters = body.get("characters") or []
    if characters:
        for page_index, page_items in enumerate(
            _paginate_items(
                characters[:8], max_items=2, max_chars=400,
                text_getter=lambda item: " ".join(str(item.get(key) or "") for key in ("name", "role", "stance")),
            ),
            1,
        ):
            slides.append({
                "kind": "characters", "section": "案例情境",
                "title": "谁在影响决策？" if page_index == 1 else f"谁在影响决策？ · {page_index}",
                "items": page_items,
                "notes": _speaker_note("关键角色", "识别正式权力、专业判断和一线信息之间的张力。", ["请学生先复述立场，不急于评价人物。"], 4),
            })

    narrative_pages = _content_pages(body.get("narrative", ""), density, DENSITY[density]["narrative_slides"])
    phase_names = ["情境建立", "冲突浮现", "证据交锋", "约束升级", "决策临界", "影响扩散", "方案博弈", "行动窗口"]
    slides.append(_section_slide("追踪冲突与证据", "02", "把观点拆成事实、假设与利益"))
    for index, items in enumerate(narrative_pages, 1):
        phase = phase_names[min(index - 1, len(phase_names) - 1)]
        slides.append({
            "kind": "story", "section": "案例进程", "title": f"{phase}", "items": items,
            "chapter": f"{index:02d}", "emphasis": index % 2, "visual_mode": mode == "visual",
            "notes": _speaker_note(phase, "按事件推进呈现信息，保持决策悬念。", ["追问：这里出现了什么新证据？", "追问：谁承担了风险？"], 3),
        })

    slides.extend(visual_pages[1:])

    if body.get("decision_point"):
        slides.append(_section_slide("站到决策者的位置", "03", "没有完美方案，只有可解释的取舍"))
        for page_index, page_items in enumerate(_decision_pages(body.get("decision_point"), density), 1):
            slides.append({
                "kind": "decision", "section": "课堂决策",
                "title": "关键决策点" if page_index == 1 else f"关键决策点 · {page_index}",
                "items": page_items,
                "criteria": ["价值与目标", "证据充分性", "风险可控性", "执行可行性"],
                "notes": _speaker_note("关键决策点", "冻结案例信息，让学生独立形成初步立场。", ["先个人写下选择与一条证据，再进入小组讨论。"], 5),
            })

    slides.append(_section_slide("研讨、辩论与复盘", "04", "用证据挑战直觉，用反馈修正方案"))
    question_pages = _paginate_items(
        questions,
        max_items=2,
        max_chars={"concise": 250, "standard": 290, "detailed": 320}[density],
        text_getter=lambda item: f"{item.get('question', '')} {item.get('teaching_intent', '') if audience == 'teacher' else ''}",
    )
    question_offset = 0
    for page_index, questions_page in enumerate(question_pages, 1):
        page_questions = [
            {**question, "_number": question_offset + offset + 1}
            for offset, question in enumerate(questions_page)
        ]
        question_offset += len(questions_page)
        slides.append({
            "kind": "questions", "section": "课堂研讨",
            "title": "课堂讨论" if page_index == 1 else f"课堂讨论 · {page_index}",
            "items": page_questions,
            "show_intent": audience == "teacher",
            "activity": "小组研讨 6 分钟 → 代表陈述 2 分钟 → 交叉质询" if mode == "workshop" else "先独立判断，再用案例证据回应",
            "notes": _speaker_note("课堂讨论", "推动学生从事实识别逐步走向评价与创造。", [*[_clean_text(q.get("teaching_intent")) for q in page_questions]], 8 if mode == "workshop" else 6),
        })

    if audience == "teacher":
        slides.append(_section_slide("教师实施与评价", "05", "把精彩讨论收束为可迁移的能力"))
        if guide.get("teaching_flow"):
            slides.append({
                "kind": "flow", "section": "教学实施", "title": "一堂课如何落地",
                "items": _split_sentences(str(guide.get("teaching_flow", "")).replace("→", "。")), "teacher_only": True,
                "notes": _speaker_note("授课流程", "帮助教师控制课堂节奏。", [guide.get("teaching_flow", "")], 2),
            })
        if guide.get("key_points") or guide.get("common_misconceptions"):
            key_pages = _paginate_items(
                guide.get("key_points") or [], max_items=2, max_chars=300, text_getter=lambda item: item,
            )
            misconception_pages = _paginate_items(
                guide.get("common_misconceptions") or [], max_items=2, max_chars=300, text_getter=lambda item: item,
            )
            tips_page_count = max(len(key_pages), len(misconception_pages))
            for page_index in range(tips_page_count):
                page_keys = key_pages[page_index] if page_index < len(key_pages) else []
                page_misconceptions = misconception_pages[page_index] if page_index < len(misconception_pages) else []
                slides.append({
                    "kind": "teaching_tips", "section": "教学实施",
                    "title": "要点与误区：教师观察清单" if page_index == 0 else f"教师观察清单 · {page_index + 1}",
                    "key_points": page_keys, "misconceptions": page_misconceptions,
                    "teacher_only": True,
                    "notes": _speaker_note("教学提示", "在讨论过程中观察学生是否抓住机制而非表象。", [*[f"要点：{x}" for x in page_keys], *[f"误区：{x}" for x in page_misconceptions]], 3),
                })
        matrix = package.get("alignment_matrix") or []
        if matrix:
            for page_index, page_items in enumerate(
                _paginate_items(
                    matrix[:8], max_items=3, max_chars=620,
                    text_getter=lambda item: " ".join(str(item.get(key) or "") for key in ("objective_id", "case_section", "activity", "assessment")),
                ),
                1,
            ):
                slides.append({
                    "kind": "alignment", "section": "教学实施",
                    "title": "目标—活动—评价闭环" if page_index == 1 else f"目标—活动—评价闭环 · {page_index}",
                    "items": page_items, "teacher_only": True,
                    "notes": _speaker_note("教学目标对齐", "确认每项学习目标都有活动与评价证据。", [], 2),
                })

    evidence_sources = package.get("evidence_sources") or []
    for index in range(0, len(evidence_sources), 5):
        page_sources = evidence_sources[index : index + 5]
        slides.append({
            "kind": "sources", "section": "事实核验", "title": "案例事实从哪里来" if index == 0 else "案例事实从哪里来 · 续",
            "items": page_sources,
            "notes": _speaker_note(
                "事实来源", "说明正文引用标记与权威来源的对应关系，明确事实、推断和教学任务的边界。",
                [f"[{item.get('id', '')}] {item.get('title', '')}" for item in page_sources], 2,
            ),
        })

    video_resources = package.get("video_resources") or []
    for index in range(0, len(video_resources), 5):
        page_videos = video_resources[index : index + 5]
        slides.append({
            "kind": "videos", "section": "拓展资源", "title": "精选课堂视频" if index == 0 else "精选课堂视频 · 续",
            "items": page_videos,
            "notes": _speaker_note("视频资源", "按教学目的选择片段，不建议课堂连续播放全部视频。", [item.get("usage", "") for item in page_videos], 2),
        })

    slides.append({
        "kind": "closing", "title": "回到决策现场",
        "content": "请基于案例证据形成判断，说明取舍依据、行动路径与风险应对。",
        "notes": _speaker_note("课堂收束", "让学生用一句话更新自己的决策，并说明改变依据。", ["可布置一页决策备忘录作为课后作业。"], 3),
    })
    return {
        "title": title, "theme": theme, "density": density, "audience": audience,
        "mode": mode, "include_speaker_notes": include_speaker_notes, "slides": slides,
        "design_metrics": _design_metrics(slides),
    }


def outline_preview(outline: dict[str, Any]) -> dict[str, Any]:
    preview = []
    for index, slide in enumerate(outline["slides"], 1):
        items = slide.get("items") or []
        if items and isinstance(items[0], dict):
            summary = "；".join(_clean_text(item.get("question") or item.get("title") or item.get("name") or item.get("objective_id") or item.get("value") or item.get("label")) for item in items[:2])
        elif items:
            summary = "；".join(_clean_text(item) for item in items[:2])
        else:
            summary = _clean_text(slide.get("content") or slide.get("subtitle") or slide.get("caption"))
        preview.append({
            "index": index, "kind": slide["kind"], "title": slide["title"],
            "summary": summary[:120], "teacher_only": bool(slide.get("teacher_only")),
        })
    return {
        **{key: outline[key] for key in ("title", "theme", "density", "audience", "mode", "include_speaker_notes")},
        "slide_count": len(preview), "slides": preview,
        "design_metrics": outline.get("design_metrics") or _design_metrics(outline["slides"]),
    }


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _shape(slide, x, y, w, h, fill: str, radius: bool = True, line: str | None = None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line or fill)
    return shape


def _textbox(slide, text: str, x: float, y: float, w: float, h: float, *, size: int = 20,
             color: str = "172033", bold: bool = False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, font: str = "Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    # WPS、PowerPoint 与 LibreOffice 对“溢出时自动缩小”的实现不一致，
    # 甚至可能反向扩张文本框。内容已在大纲阶段分页，这里固定文本框边界。
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = Inches(0.05)
    frame.margin_top = frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _fit_font_size(text: Any, w: float, h: float, maximum: int, minimum: int = 10) -> int:
    """按中文全角字符估算可读字号；分页负责控量，此处只做最后适配。"""
    length = max(len(_clean_text(text)), 1)
    for size in range(maximum, minimum - 1, -1):
        capacity = w * h * 72 * 72 / max(size * size * 1.35, 1)
        if capacity >= length:
            return size
    return minimum


def _add_notes(slide, text: str | None) -> None:
    if not text:
        return
    frame = slide.notes_slide.notes_text_frame
    frame.text = text


def _add_image_cover(slide, image_path: str | Path, x: float, y: float, w: float, h: float):
    picture = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    width_px, height_px = picture.image.size
    image_ratio = width_px / max(height_px, 1)
    target_ratio = w / max(h, 0.01)
    if image_ratio > target_ratio:
        crop = (1 - target_ratio / image_ratio) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < target_ratio:
        crop = (1 - image_ratio / target_ratio) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def _validate_slide_bounds(prs: Presentation) -> None:
    """发布门禁：任何可见图形都不得使用负坐标或超出 16:9 画布。"""
    errors: list[str] = []
    tolerance = Inches(0.01)
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape_index, shape in enumerate(slide.shapes, 1):
            if shape.left < -tolerance or shape.top < -tolerance:
                errors.append(f"第{slide_index}页图形{shape_index}使用负坐标")
            if shape.left + shape.width > prs.slide_width + tolerance:
                errors.append(f"第{slide_index}页图形{shape_index}超出右边界")
            if shape.top + shape.height > prs.slide_height + tolerance:
                errors.append(f"第{slide_index}页图形{shape_index}超出下边界")
    if errors:
        raise ValueError("PPT 版式边界校验失败：" + "；".join(errors[:8]))


def _base_slide(prs: Presentation, palette: dict[str, str], title: str, section: str | None, page: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _rgb(palette["background"])
    style = palette.get("style", "academic")
    if style == "editorial":
        _shape(slide, 0.65, 0.38, 1.15, 0.07, palette["accent"], radius=False)
        _textbox(slide, f"{page:02d}", 11.55, 0.25, 0.9, 0.55, size=30, color=palette["soft"], bold=True, align=PP_ALIGN.RIGHT)
    elif style == "executive":
        _shape(slide, 0, 0, 13.333, 1.35, palette["primary"], radius=False)
        _shape(slide, 0.65, 1.31, 2.15, 0.04, palette["accent"], radius=False)
    elif style == "vibrant":
        _shape(slide, 0, 0, 0.16, 7.5, palette["accent2"], radius=False)
        _shape(slide, 11.98, 0.0, 1.353, 0.95, palette["soft"], radius=True)
        _shape(slide, 12.47, 0.15, 0.46, 0.46, palette["accent2"], radius=True)
    else:
        _shape(slide, 0, 0, 0.13, 7.5, palette["accent"], radius=False)
        if style == "academic":
            _shape(slide, 11.98, 0.0, 1.353, 0.95, palette["soft"], radius=True)
            _shape(slide, 12.47, 0.15, 0.46, 0.46, palette["accent2"], radius=True)
    header_color = "FFFFFF" if style == "executive" else palette["primary"]
    section_color = palette["accent"] if style != "executive" else palette["soft"]
    if section:
        _textbox(slide, section.upper(), 0.65, 0.34, 3.5, 0.3, size=9, color=section_color, bold=True)
    title_size = _fit_font_size(title, 11.0, 0.55, 25, 19)
    _textbox(slide, title, 0.65, 0.72, 11.0, 0.58, size=title_size, color=header_color, bold=True)
    _shape(slide, 0.65, 7.08, 12.0, 0.02, "D9E2E7", radius=False)
    _shape(slide, 0.65, 7.08, 12.0 * page / max(total, 1), 0.02, palette["accent"], radius=False)
    _textbox(slide, "知案 · AI 教学案例课件", 0.65, 7.16, 4.0, 0.2, size=8, color=palette["muted"])
    _textbox(slide, f"{page:02d}", 11.9, 7.13, 0.75, 0.23, size=9, color=palette["muted"], bold=True, align=PP_ALIGN.RIGHT)
    return slide


def _render_story(slide, spec: dict[str, Any], palette: dict[str, str]) -> None:
    items = spec.get("items") or []
    emphasis = int(spec.get("emphasis") or 0)
    if spec.get("visual_mode"):
        _shape(slide, 0.72, 1.5, 3.0, 4.85, palette["primary"])
        _textbox(slide, str(spec.get("chapter") or "STORY"), 1.05, 1.9, 2.25, 0.8, size=38, color=palette["accent2"], bold=True)
        _textbox(slide, "观察 · 判断 · 追问", 1.05, 5.55, 2.25, 0.35, size=10, color="DCE7EC", bold=True)
        if items:
            _textbox(slide, items[0], 4.15, 1.7, 7.95, 2.0, size=25, color=palette["primary"], bold=True, valign=MSO_ANCHOR.MIDDLE)
        for index, item in enumerate((items[1:] or ["请结合案例原文识别新信息与约束。"] )[:2]):
            y = 4.08 + index * 1.05
            _shape(slide, 4.18, y + 0.08, 0.22, 0.22, palette["accent2"])
            _textbox(slide, item, 4.62, y, 7.3, 0.72, size=15, color=palette["text"], valign=MSO_ANCHOR.MIDDLE)
        return
    _textbox(slide, str(spec.get("chapter") or "STORY"), 0.72, 1.52, 1.1, 0.85, size=34, color=palette["soft"], bold=True)
    if items:
        _textbox(slide, items[0], 1.83, 1.56, 10.25, 1.18, size=23, color=palette["primary"], bold=True, valign=MSO_ANCHOR.MIDDLE)
    rest = items[1:] or ["请结合案例原文，识别这一阶段出现的新信息与约束。"]
    card_fill = palette["soft"] if emphasis else palette["surface"]
    for index, item in enumerate(rest[:4]):
        col = index % 2
        row = index // 2
        x, y = 0.75 + col * 6.0, 3.05 + row * 1.56
        _shape(slide, x, y, 5.65, 1.28, card_fill, line="D9E2E7")
        _shape(slide, x + 0.22, y + 0.25, 0.18, 0.18, palette["accent2"], radius=True)
        _textbox(slide, item, x + 0.55, y + 0.18, 4.78, 0.86, size=15, color=palette["text"], valign=MSO_ANCHOR.MIDDLE)


def _render_overview(slide, spec: dict[str, Any], palette: dict[str, str]) -> None:
    items = spec.get("items") or []
    for index, item in enumerate(items[:4]):
        x = 0.72 + (index % 2) * 4.45
        y = 1.58 + (index // 2) * 1.55
        _shape(slide, x, y, 4.1, 1.22, palette["surface"], line="D9E2E7")
        _textbox(slide, item.get("label", ""), x + 0.25, y + 0.2, 0.85, 0.28, size=10, color=palette["muted"], bold=True)
        _textbox(slide, item.get("value", ""), x + 0.25, y + 0.55, 3.55, 0.4, size=18, color=palette["primary"], bold=True)
    _shape(slide, 9.78, 1.58, 2.72, 4.38, palette["primary"])
    for index, stat in enumerate(spec.get("stats") or []):
        y = 1.94 + index * 1.25
        _textbox(slide, str(stat.get("value", 0)), 10.05, y, 0.8, 0.55, size=29, color="FFFFFF", bold=True)
        _textbox(slide, stat.get("label", ""), 10.88, y + 0.12, 1.28, 0.3, size=11, color="DCE7EC", bold=True)
    _textbox(slide, "CASE\nAT A GLANCE", 0.74, 5.3, 5.0, 0.78, size=12, color=palette["accent"], bold=True)


def _render_objectives(slide, spec: dict[str, Any], palette: dict[str, str]) -> None:
    items = spec.get("items") or []
    levels = spec.get("levels") or []
    count = max(len(items), 1)
    gap = 0.28
    card_w = min(5.65, (11.6 - gap * (count - 1)) / count)
    start_number = int(spec.get("start_number") or 1)
    for index, item in enumerate(items[:4]):
        x = 0.72 + index * (card_w + gap)
        y = 1.72
        h = 4.55
        _shape(slide, x, y, card_w, h, palette["surface"], line="D9E2E7")
        _shape(slide, x, y, card_w, 0.13, [palette["accent"], palette["accent2"], palette["success"], palette["warning"]][index % 4], radius=False)
        _textbox(slide, f"{start_number + index:02d}", x + 0.25, y + 0.33, 0.7, 0.48, size=24, color=palette["soft"], bold=True)
        if index < len(levels):
            _shape(slide, x + card_w - 1.0, y + 0.38, 0.72, 0.34, palette["soft"])
            _textbox(slide, levels[index], x + card_w - 0.94, y + 0.45, 0.6, 0.18, size=9, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)
        font_size = _fit_font_size(item, card_w - 0.5, 2.55, 18, 12)
        _textbox(slide, item, x + 0.25, y + 1.05, card_w - 0.5, 2.65, size=font_size, color=palette["text"], bold=True, valign=MSO_ANCHOR.MIDDLE)
        _textbox(slide, "学习产出", x + 0.25, y + h - 0.62, card_w - 0.5, 0.22, size=9, color=palette["muted"], bold=True)


def _render_bullets(slide, items: list[str], palette: dict[str, str], *, start_y: float = 1.62):
    count = max(len(items), 1)
    available = 5.05
    row_h = min(1.15, available / count)
    font_size = 20 if count <= 4 else 17
    for index, item in enumerate(items):
        y = start_y + index * row_h
        _shape(slide, 0.72, y + 0.06, 0.38, 0.38, palette["soft"], radius=True)
        _textbox(slide, f"{index + 1}", 0.72, y + 0.075, 0.38, 0.25, size=10, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, _clean_text(item), 1.28, y, 11.1, row_h - 0.05, size=font_size, color=palette["text"], valign=MSO_ANCHOR.MIDDLE)


def _render_method_map(slide, spec: dict[str, Any], palette: dict[str, str]) -> None:
    items = spec.get("items") or []
    colors = [palette["accent"], palette["accent2"], palette["success"], palette["warning"], palette["primary"], palette["danger"]]
    for index, item in enumerate(items[:6]):
        col, row = index % 3, index // 3
        x, y = 0.74 + col * 4.03, 1.65 + row * 1.65
        _shape(slide, x + 0.07, y + 0.08, 3.65, 1.25, "D9E2E7")
        _shape(slide, x, y, 3.65, 1.25, palette["surface"], line="D9E2E7")
        _shape(slide, x + 0.22, y + 0.28, 0.55, 0.55, colors[index])
        _textbox(slide, f"{index + 1:02d}", x + 0.29, y + 0.43, 0.4, 0.18, size=9, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, item, x + 0.98, y + 0.25, 2.35, 0.66, size=15, color=palette["text"], bold=True, valign=MSO_ANCHOR.MIDDLE)
        if index < len(items[:6]) - 1:
            connector = "↓" if index == 2 else "→"
            cx = x + (1.54 if index == 2 else 3.66)
            cy = y + (1.15 if index == 2 else 0.42)
            _textbox(slide, connector, cx, cy, 0.35, 0.35, size=14, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)
    problem = _clean_text(spec.get("problem"))
    if problem:
        _shape(slide, 0.74, 5.13, 11.72, 0.92, palette["primary"])
        _textbox(slide, "教学难点", 1.02, 5.36, 0.92, 0.25, size=10, color=palette["accent2"], bold=True)
        _textbox(slide, problem[:120], 2.08, 5.26, 9.96, 0.48, size=12, color="FFFFFF", bold=True, valign=MSO_ANCHOR.MIDDLE)


def _render_ideology(slide, spec: dict[str, Any], palette: dict[str, str]) -> None:
    _shape(slide, 0.72, 1.55, 4.1, 4.78, palette["primary"])
    _textbox(slide, "VALUE\nCOMPASS", 1.08, 1.95, 2.7, 0.85, size=12, color=palette["accent2"], bold=True)
    _textbox(slide, spec.get("figure") or "实践人物", 1.08, 3.0, 3.05, 0.7, size=31, color="FFFFFF", bold=True)
    _textbox(slide, "用真实贡献呈现价值，不虚构对白与名言", 1.08, 4.72, 3.12, 0.75, size=13, color="DCE7EC", bold=True)
    themes = spec.get("themes") or []
    for index, theme in enumerate(themes[:6]):
        col, row = index % 2, index // 2
        x, y = 5.25 + col * 3.42, 1.72 + row * 0.9
        _shape(slide, x, y, 3.08, 0.65, palette["soft"], line="D9E2E7")
        _textbox(slide, theme, x + 0.15, y + 0.18, 2.78, 0.25, size=13, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)
    _shape(slide, 5.25, 4.68, 7.18, 1.48, palette["surface"], line="D9E2E7")
    _textbox(slide, "如何融入案例", 5.55, 4.95, 1.2, 0.28, size=10, color=palette["muted"], bold=True)
    _textbox(slide, spec.get("implementation", ""), 6.76, 4.84, 5.22, 0.72, size=13, color=palette["text"], bold=True, valign=MSO_ANCHOR.MIDDLE)


def _render_visual_gallery(slide, spec: dict[str, Any], palette: dict[str, str]) -> None:
    items = spec.get("items") or []
    count = min(len(items), 4)
    layouts = {
        1: [(0.72, 1.5, 11.82, 4.95)],
        2: [(0.72, 1.5, 7.22, 4.95), (8.16, 1.5, 4.38, 4.95)],
        3: [(0.72, 1.5, 7.22, 4.95), (8.16, 1.5, 4.38, 2.36), (8.16, 4.09, 4.38, 2.36)],
        4: [(0.72, 1.5, 5.76, 2.36), (6.74, 1.5, 5.8, 2.36), (0.72, 4.09, 5.76, 2.36), (6.74, 4.09, 5.8, 2.36)],
    }
    for index, item in enumerate(items[:4]):
        x, y, w, h = layouts[count][index]
        try:
            _add_image_cover(slide, get_cached_material_image(item["id"]), x, y, w, h)
        except Exception:
            _shape(slide, x, y, w, h, palette["soft"], line="D9E2E7")
            _textbox(slide, "图片暂未缓存", x + 0.2, y + h / 2 - 0.2, w - 0.4, 0.35, size=11, color=palette["muted"], align=PP_ALIGN.CENTER)
        _shape(slide, x, y + h - 0.62, w, 0.62, palette["primary"], radius=False)
        caption = _clean_text(item.get("caption") or item.get("title"))
        _textbox(slide, caption[:46] + ("…" if len(caption) > 46 else ""), x + 0.18, y + h - 0.5, w - 0.36, 0.22, size=9 if count >= 3 else 11, color="FFFFFF", bold=True)
        source = _textbox(slide, f"{item.get('source_org', '')}  ↗", x + 0.18, y + h - 0.25, w - 0.36, 0.16, size=7, color="DCE7EC", align=PP_ALIGN.RIGHT)
        if item.get("source_page_url"):
            source.click_action.hyperlink.address = item["source_page_url"]


def export_pptx(package: dict[str, Any], title: str, version: int = 1, options: dict[str, Any] | None = None) -> str:
    outline = build_ppt_outline(package, options)
    palette = THEMES[outline["theme"]]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = outline["title"]
    prs.core_properties.subject = "教学案例课堂课件"
    prs.core_properties.author = "知案 · AI 教学案例工作台"

    total_slides = len(outline["slides"])
    for page, spec in enumerate(outline["slides"], 1):
        kind = spec["kind"]
        if kind == "cover":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            fill = slide.background.fill
            fill.solid(); fill.fore_color.rgb = _rgb(palette["primary"])
            if spec.get("asset_id"):
                try:
                    _add_image_cover(slide, get_cached_material_image(spec["asset_id"]), 7.75, 0, 5.58, 7.5)
                    _shape(slide, 7.5, 0, 0.3, 7.5, palette["accent"], radius=False)
                except Exception:
                    pass
            _shape(slide, 0.72, 0.72, 1.15, 0.1, palette["accent2"], radius=False)
            _textbox(slide, spec.get("eyebrow", "教学案例课堂课件"), 0.72, 1.05, 5.9, 0.35, size=11, color="DCE7EC", bold=True)
            _textbox(slide, spec["title"], 0.72, 1.62, 6.45, 2.75, size=31, color="FFFFFF", bold=True, valign=MSO_ANCHOR.MIDDLE)
            _textbox(slide, spec.get("subtitle", ""), 0.74, 4.65, 6.15, 0.55, size=14, color="DCE7EC")
            _shape(slide, 0.72, 5.55, 6.2, 0.02, palette["accent2"], radius=False)
            _textbox(slide, f"V{version}  ·  {datetime.now():%Y-%m-%d}", 0.74, 5.82, 2.8, 0.35, size=10, color="B9CDD6")
            _textbox(slide, "知案 · AI 教学案例工作台", 3.4, 5.82, 3.52, 0.35, size=10, color="B9CDD6", align=PP_ALIGN.RIGHT)
            if outline["include_speaker_notes"]:
                _add_notes(slide, spec.get("notes"))
            continue

        if kind == "section":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = _rgb(palette["primary"])
            _textbox(slide, spec.get("number", ""), 0.72, 0.68, 2.2, 1.2, size=56, color=palette["accent"], bold=True)
            _shape(slide, 0.76, 2.12, 1.3, 0.09, palette["accent2"], radius=False)
            _textbox(slide, spec["title"], 0.72, 2.55, 11.6, 1.35, size=36, color="FFFFFF", bold=True, valign=MSO_ANCHOR.MIDDLE)
            _textbox(slide, spec.get("kicker", ""), 0.75, 4.25, 9.8, 0.58, size=17, color="DCE7EC")
            _textbox(slide, "SECTION", 10.0, 6.45, 2.5, 0.28, size=10, color="B9CDD6", bold=True, align=PP_ALIGN.RIGHT)
            if outline["include_speaker_notes"]:
                _add_notes(slide, spec.get("notes"))
            continue

        slide = _base_slide(prs, palette, spec["title"], spec.get("section"), page, total_slides)
        if spec.get("teacher_only"):
            _shape(slide, 10.85, 0.38, 1.55, 0.38, palette["soft"], radius=True)
            _textbox(slide, "教师版专属", 10.9, 0.45, 1.42, 0.2, size=9, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)

        if kind == "objectives":
            _render_objectives(slide, spec, palette)
        elif kind == "overview":
            _render_overview(slide, spec, palette)
        elif kind == "story":
            _render_story(slide, spec, palette)
        elif kind == "method_map":
            _render_method_map(slide, spec, palette)
        elif kind == "ideology":
            _render_ideology(slide, spec, palette)
        elif kind == "visual_gallery":
            _render_visual_gallery(slide, spec, palette)
        elif kind == "agenda":
            items = spec.get("items") or []
            _shape(slide, 1.18, 3.17, 10.72, 0.06, "D9E2E7", radius=False)
            for index, item in enumerate(items[:6]):
                x = 1.0 + index * (10.85 / max(len(items) - 1, 1))
                color = [palette["accent"], palette["accent2"], palette["success"], palette["warning"], palette["primary"], palette["danger"]][index]
                _shape(slide, x, 2.76, 0.75, 0.75, color)
                _textbox(slide, str(index + 1), x + 0.08, 2.99, 0.59, 0.24, size=12, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
                _textbox(slide, item, x - 0.42, 3.76, 1.62, 0.65, size=13, color=palette["text"], bold=True, align=PP_ALIGN.CENTER)
            _textbox(slide, "不是听完一个故事，而是完成一次判断练习", 1.05, 5.35, 11.2, 0.42, size=16, color=palette["muted"], align=PP_ALIGN.CENTER)
        elif kind == "bullets":
            _render_bullets(slide, spec.get("items") or [], palette)
        elif kind == "visual":
            try:
                image_path = get_cached_material_image(spec["asset_id"])
                _add_image_cover(slide, image_path, 0.72, 1.48, 8.15, 4.95)
            except Exception:
                _shape(slide, 0.72, 1.48, 8.05, 4.95, palette["soft"], line="D9E2E7")
                _textbox(slide, "官方图片暂时无法获取\n请通过来源页面查看", 1.2, 3.05, 7.0, 0.9, size=18, color=palette["muted"], align=PP_ALIGN.CENTER)
            _shape(slide, 9.08, 1.48, 3.45, 4.95, palette["surface"], line="D9E2E7")
            _textbox(slide, spec.get("caption", ""), 9.38, 1.86, 2.85, 1.6, size=16, color=palette["text"], bold=True)
            _textbox(slide, f"来源机构\n{spec.get('source_org', '')}", 9.38, 3.72, 2.85, 0.72, size=11, color=palette["muted"])
            _textbox(slide, f"摄影\n{spec.get('photographer') or '未标注'}", 9.38, 4.62, 2.85, 0.65, size=11, color=palette["muted"])
            source_box = _textbox(slide, "查看官方原始页面 ↗", 9.38, 5.38, 2.85, 0.28, size=9, color=palette["accent"], bold=True)
            if spec.get("source_page_url"):
                source_box.click_action.hyperlink.address = spec["source_page_url"]
            _textbox(slide, "教学引用 · 外部分发前确认授权", 9.38, 5.78, 2.85, 0.3, size=8, color=palette["muted"])
        elif kind == "characters":
            items = spec.get("items") or []
            count = max(len(items), 1)
            gap = 0.24
            card_w, card_h = (11.82 - gap * (count - 1)) / count, 4.7
            for index, item in enumerate(items):
                x = 0.72 + index * (card_w + gap)
                y = 1.55
                _shape(slide, x, y, card_w, card_h, palette["surface"], line="D9E2E7")
                _shape(slide, x, y, 0.12, card_h, palette["accent"], radius=False)
                _textbox(slide, _clean_text(item.get("name")), x + 0.3, y + 0.3, card_w - 0.6, 0.55, size=18, color=palette["primary"], bold=True)
                _textbox(slide, _clean_text(item.get("role")), x + 0.3, y + 1.02, card_w - 0.6, 0.65, size=10, color=palette["muted"], bold=True)
                stance = _clean_text(item.get("stance"))
                stance_size = _fit_font_size(stance, card_w - 0.6, 2.25, 15, 10)
                _textbox(slide, stance, x + 0.3, y + 1.78, card_w - 0.6, 2.25, size=stance_size, color=palette["text"], valign=MSO_ANCHOR.TOP)
        elif kind == "decision":
            _shape(slide, 0.72, 1.55, 7.55, 4.55, palette["primary"])
            _textbox(slide, "DECISION MOMENT", 1.05, 1.93, 2.6, 0.28, size=10, color=palette["accent2"], bold=True)
            decision_text = "\n\n".join(f"{index + 1}. {_clean_text(item)}" for index, item in enumerate(spec.get("items") or []))
            decision_size = _fit_font_size(decision_text, 6.85, 2.75, 18, 12)
            _textbox(slide, decision_text, 1.05, 2.35, 6.85, 2.85, size=decision_size, color="FFFFFF", bold=True, valign=MSO_ANCHOR.MIDDLE)
            _textbox(slide, "先选择，再说明你愿意承担什么风险。", 1.05, 5.44, 6.65, 0.3, size=11, color="DCE7EC")
            _textbox(slide, "决策检验框架", 8.72, 1.62, 3.4, 0.35, size=12, color=palette["muted"], bold=True)
            for index, criterion in enumerate(spec.get("criteria") or []):
                y = 2.18 + index * 0.95
                _shape(slide, 8.68, y, 3.72, 0.72, palette["surface"], line="D9E2E7")
                _textbox(slide, f"{index + 1}", 8.9, y + 0.19, 0.36, 0.24, size=10, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)
                _textbox(slide, criterion, 9.43, y + 0.14, 2.55, 0.32, size=14, color=palette["text"], bold=True)
        elif kind == "questions":
            items = spec.get("items") or []
            _shape(slide, 0.72, 1.47, 11.82, 0.52, palette["soft"])
            _textbox(slide, spec.get("activity", "先独立判断，再用案例证据回应"), 0.98, 1.62, 10.9, 0.22, size=10, color=palette["accent"], bold=True)
            row_h = 1.82 if len(items) >= 2 else 3.82
            for index, item in enumerate(items):
                y = 2.15 + index * (row_h + 0.12)
                _shape(slide, 0.72, y, 11.82, row_h, palette["surface"], line="D9E2E7")
                _textbox(slide, f"Q{item.get('_number', index + 1)}", 1.0, y + 0.3, 0.65, 0.35, size=13, color=palette["accent"], bold=True)
                question_text = f"[{item.get('level', '')}] {_clean_text(item.get('question'))}"
                question_h = row_h - (0.62 if spec.get("show_intent") and item.get("teaching_intent") else 0.3)
                question_size = _fit_font_size(question_text, 10.1, question_h, 17, 12)
                _textbox(slide, question_text, 1.72, y + 0.16, 10.1, question_h, size=question_size, color=palette["text"], bold=True, valign=MSO_ANCHOR.MIDDLE)
                if spec.get("show_intent") and item.get("teaching_intent"):
                    _textbox(slide, f"教学意图 · {_clean_text(item.get('teaching_intent'))}", 1.72, y + row_h - 0.43, 9.9, 0.24, size=9, color=palette["muted"])
        elif kind == "flow":
            items = spec.get("items") or []
            count = max(len(items), 1)
            card_w = min(2.75, 11.25 / count)
            for index, item in enumerate(items[:5]):
                x = 0.72 + index * (card_w + 0.16)
                _shape(slide, x, 2.0, card_w, 3.55, palette["surface"], line="D9E2E7")
                _shape(slide, x, 2.0, card_w, 0.12, [palette["accent"], palette["accent2"], palette["success"], palette["warning"], palette["primary"]][index], radius=False)
                _textbox(slide, f"{index + 1:02d}", x + 0.22, 2.38, 0.65, 0.42, size=22, color=palette["soft"], bold=True)
                _textbox(slide, item, x + 0.22, 3.05, card_w - 0.44, 1.75, size=15, color=palette["text"], bold=True, valign=MSO_ANCHOR.MIDDLE)
                if index < len(items[:5]) - 1:
                    _textbox(slide, "→", x + card_w - 0.08, 3.38, 0.32, 0.4, size=17, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)
        elif kind == "teaching_tips":
            _shape(slide, 0.72, 1.58, 5.72, 4.82, "ECFDF5", line="A7F3D0")
            _shape(slide, 6.67, 1.58, 5.87, 4.82, "FEF2F2", line="FECACA")
            _textbox(slide, "✓ 应该抓住", 1.02, 1.93, 4.9, 0.35, size=15, color=palette["success"], bold=True)
            _textbox(slide, "! 需要警惕", 6.98, 1.93, 4.9, 0.35, size=15, color=palette["danger"], bold=True)
            key_points = spec.get("key_points") or []
            misconceptions = spec.get("misconceptions") or []
            for index, item in enumerate(key_points):
                row_h = 3.45 / max(len(key_points), 1)
                y = 2.58 + index * row_h
                text = f"{index + 1:02d}  {_clean_text(item)}"
                size = _fit_font_size(text, 4.92, row_h - 0.12, 14, 10)
                _textbox(slide, text, 1.02, y, 4.92, row_h - 0.12, size=size, color=palette["text"], bold=True, valign=MSO_ANCHOR.MIDDLE)
            for index, item in enumerate(misconceptions):
                row_h = 3.45 / max(len(misconceptions), 1)
                y = 2.58 + index * row_h
                text = f"{index + 1:02d}  {_clean_text(item)}"
                size = _fit_font_size(text, 5.02, row_h - 0.12, 14, 10)
                _textbox(slide, text, 6.98, y, 5.02, row_h - 0.12, size=size, color=palette["text"], bold=True, valign=MSO_ANCHOR.MIDDLE)
        elif kind == "alignment":
            rows = spec.get("items") or []
            headers = ["目标", "案例环节", "课堂活动", "评价方式"]
            widths = [1.4, 3.6, 3.2, 3.6]
            x_positions = [0.72]
            for width in widths[:-1]:
                x_positions.append(x_positions[-1] + width)
            for col, header in enumerate(headers):
                _shape(slide, x_positions[col], 1.58, widths[col], 0.55, palette["primary"], radius=False)
                _textbox(slide, header, x_positions[col] + 0.08, 1.73, widths[col] - 0.16, 0.23, size=10, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
            row_h = min(1.36, 4.08 / max(len(rows), 1))
            for row_index, row in enumerate(rows):
                values = [row.get(key, "") for key in ("objective_id", "case_section", "activity", "assessment")]
                y = 2.13 + row_index * row_h
                for col, value in enumerate(values):
                    fill = "FFFFFF" if row_index % 2 == 0 else palette["background"]
                    _shape(slide, x_positions[col], y, widths[col], row_h, fill, radius=False, line="D9E2E7")
                    clean_value = _clean_text(value)
                    size = _fit_font_size(clean_value, widths[col] - 0.16, row_h - 0.1, 11, 8)
                    _textbox(slide, clean_value, x_positions[col] + 0.08, y + 0.05, widths[col] - 0.16, row_h - 0.1, size=size, color=palette["text"], valign=MSO_ANCHOR.MIDDLE)
        elif kind == "sources":
            for index, source in enumerate(spec.get("items") or []):
                y = 1.48 + index * 1.04
                _shape(slide, 0.72, y, 11.82, 0.88, palette["surface"], line="D9E2E7")
                _textbox(slide, f"[{source.get('id', '')}]", 0.98, y + 0.17, 0.62, 0.26, size=11, color=palette["accent"], bold=True)
                _textbox(slide, _clean_text(source.get("title")), 1.68, y + 0.1, 7.55, 0.3, size=13, color=palette["text"], bold=True)
                _textbox(slide, _clean_text(source.get("source_org")), 9.38, y + 0.12, 2.75, 0.24, size=9, color=palette["muted"], align=PP_ALIGN.RIGHT)
                _textbox(slide, _clean_text(source.get("usage"))[:105], 1.68, y + 0.5, 9.85, 0.22, size=8, color=palette["muted"])
                link = _textbox(slide, "官网原文 ↗", 11.38, y + 0.49, 0.75, 0.2, size=7, color=palette["accent"], bold=True, align=PP_ALIGN.RIGHT)
                if source.get("source_page_url"):
                    link.click_action.hyperlink.address = source["source_page_url"]
        elif kind == "videos":
            for index, video in enumerate(spec.get("items") or []):
                y = 1.48 + index * 1.04
                _shape(slide, 0.72, y, 11.82, 0.88, palette["surface"], line="D9E2E7")
                _shape(slide, 0.98, y + 0.17, 0.5, 0.5, palette["accent"])
                _textbox(slide, "▶", 1.06, y + 0.29, 0.34, 0.18, size=9, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
                _textbox(slide, _clean_text(video.get("title")), 1.66, y + 0.09, 7.9, 0.3, size=13, color=palette["text"], bold=True)
                _textbox(slide, _clean_text(video.get("source_org")), 9.5, y + 0.11, 2.55, 0.23, size=8, color=palette["muted"], align=PP_ALIGN.RIGHT)
                _textbox(slide, _clean_text(video.get("usage"))[:105], 1.66, y + 0.49, 9.45, 0.22, size=8, color=palette["muted"])
                link = _textbox(slide, "播放 ↗", 11.28, y + 0.49, 0.75, 0.2, size=7, color=palette["accent"], bold=True, align=PP_ALIGN.RIGHT)
                if video.get("video_url"):
                    link.click_action.hyperlink.address = video["video_url"]
        elif kind == "closing":
            _shape(slide, 0.72, 1.8, 11.82, 3.7, palette["primary"], radius=True)
            _textbox(slide, spec.get("content", ""), 1.2, 2.34, 10.85, 1.55, size=27, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            _textbox(slide, "EVIDENCE · TRADE-OFF · ACTION", 2.1, 4.35, 9.0, 0.35, size=11, color="DCE7EC", bold=True, align=PP_ALIGN.CENTER)

        if outline["include_speaker_notes"]:
            _add_notes(slide, spec.get("notes"))

    clean_title = re.sub(r"[\\/:*?\"<>|\x00-\x1f]", "_", outline["title"]).strip(" ._")[:80] or title
    theme_name = THEME_NAMES.get(outline["theme"], outline["theme"])
    filename = f"{clean_title}_教学案例课件_{theme_name}_V{version}.pptx"
    export_dir = Path(settings.export_dir)
    export_dir.mkdir(parents=True, exist_ok=True)
    path = export_dir / filename
    _validate_slide_bounds(prs)
    prs.save(str(path))
    return str(path)
